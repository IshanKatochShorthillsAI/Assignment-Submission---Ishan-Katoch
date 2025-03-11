import base64
import re
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL


def merge_lines(raw_text):
    lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
    merged = []
    buffer = ""
    for line in lines:
        if not buffer:
            buffer = line
        else:
            if re.search(r'[.!?:"\']\s*$', buffer):
                merged.append(buffer)
                buffer = line
            else:
                buffer += " " + line
    if buffer:
        merged.append(buffer)
    return "\n\n".join(merged)


def infer_font_style(font_name):
    if not font_name:
        return False, False
    bold = "Bold" in font_name
    italic = ("Italic" in font_name) or ("Oblique" in font_name)
    return bold, italic


def extract_images_from_rels(document):
    """
    Extract images from a PPTX file by iterating over each slide's relationships.
    Returns a list of image dictionaries.
    """
    images = []
    for idx, slide in enumerate(document.slides):
        # Use the slide index (idx+1) for numbering.
        for rel in slide.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_part = rel.target_part
                    image_blob = image_part.blob
                    encoded_blob = base64.b64encode(image_blob).decode("utf-8")
                    images.append(
                        {
                            "slide": idx + 1,
                            "rId": rel.rId,
                            "format": getattr(image_part, "ext", "unknown"),
                            "blob": encoded_blob,
                        }
                    )
                except Exception:
                    continue
    return images


class DataExtractor:
    def __init__(self, loader):
        self.loader = loader

    def extract_text(self, filepath: str):
        document = self.loader.load_file(filepath)
        text_data = []
        try:
            import fitz  # PyMuPDF

            is_mupdf = hasattr(document, "page_count")
        except ImportError:
            is_mupdf = False

        if is_mupdf:
            doc_meta = document.metadata
            text_data.append({"document_metadata": doc_meta})
            for i in range(document.page_count):
                page = document[i]
                page_dict = page.get_text("dict")
                for block in page_dict.get("blocks", []):
                    if block.get("type") == 0:
                        for line in block.get("lines", []):
                            spans = line.get("spans", [])
                            line_text = " ".join(
                                [span.get("text", "").strip() for span in spans]
                            )
                            line_text = re.sub(r"\s+", " ", line_text).strip()
                            if not line_text:
                                continue
                            spans_meta = []
                            for span in spans:
                                bold, italic = infer_font_style(span.get("font", ""))
                                spans_meta.append(
                                    {
                                        "text": span.get("text", ""),
                                        "font": span.get("font", ""),
                                        "size": span.get("size", 0),
                                        "bbox": span.get("bbox", []),
                                        "flags": span.get("flags"),
                                        "origin": span.get("origin"),
                                        "color": span.get("color"),
                                        "bold": bold,
                                        "italic": italic,
                                        "additional": {
                                            k: v
                                            for k, v in span.items()
                                            if k
                                            not in [
                                                "text",
                                                "font",
                                                "size",
                                                "bbox",
                                                "flags",
                                                "origin",
                                                "color",
                                            ]
                                        },
                                    }
                                )
                            text_data.append(
                                {
                                    "page": page.number + 1,
                                    "rotation": page.rotation,
                                    "dimensions": [
                                        page.rect.x0,
                                        page.rect.y0,
                                        page.rect.x1,
                                        page.rect.y1,
                                    ],
                                    "line_text": line_text,
                                    "spans": spans_meta,
                                }
                            )
            return text_data
        elif hasattr(document, "pages"):
            try:
                import pdfplumber
            except ImportError:
                pdfplumber = None
            if pdfplumber:
                try:
                    with pdfplumber.open(filepath) as pdf:
                        for i, page in enumerate(pdf.pages):
                            text = page.extract_text(layout=True)
                            text = re.sub(r"\s+", " ", text).strip() if text else ""
                            text_data.append({"page": i + 1, "text": text})
                except Exception as e:
                    for i, page in enumerate(document.pages):
                        try:
                            raw_text = (
                                page.extract_text()
                                if hasattr(page, "extract_text")
                                else ""
                            )
                            text = merge_lines(raw_text)
                        except Exception as ex:
                            text = f"Error extracting text: {str(ex)}"
                        text_data.append({"page": i + 1, "text": text})
            else:
                for i, page in enumerate(document.pages):
                    try:
                        raw_text = (
                            page.extract_text() if hasattr(page, "extract_text") else ""
                        )
                        text = merge_lines(raw_text)
                    except Exception as e:
                        text = f"Error extracting text: {str(e)}"
                    text_data.append({"page": i + 1, "text": text})
            return text_data
        elif hasattr(document, "paragraphs"):
            for para in document.paragraphs:
                para_data = {
                    "text": para.text.strip() if para.text else "",
                    "style": str(para.style) if hasattr(para, "style") else None,
                    "alignment": para.alignment if hasattr(para, "alignment") else None,
                    "runs": [],
                }
                if hasattr(para, "runs"):
                    for run in para.runs:
                        run_data = {
                            "text": run.text,
                            "font": (
                                run.font.name
                                if run.font and hasattr(run.font, "name")
                                else "DefaultFont"
                            ),
                            "size": (
                                run.font.size.pt if run.font and run.font.size else 11
                            ),
                            "bold": run.bold if run.bold is not None else False,
                            "italic": run.italic if run.italic is not None else False,
                            "underline": (
                                run.underline if run.underline is not None else False
                            ),
                            "color": (
                                run.font.color.rgb
                                if run.font
                                and run.font.color
                                and hasattr(run.font.color, "rgb")
                                else "000000"
                            ),
                            "highlight": (
                                run.font.highlight_color
                                if run.font and hasattr(run.font, "highlight_color")
                                else None
                            ),
                        }
                        para_data["runs"].append(run_data)
                text_data.append(para_data)
            return text_data
        elif hasattr(document, "slides"):
            for i, slide in enumerate(document.slides):
                slide_items = []
                for shape in slide.shapes:
                    shape_data = {
                        "shape_id": getattr(shape, "shape_id", None),
                        "left": getattr(shape, "left", None),
                        "top": getattr(shape, "top", None),
                        "width": getattr(shape, "width", None),
                        "height": getattr(shape, "height", None),
                    }
                    if hasattr(shape, "text_frame") and shape.text_frame is not None:
                        paragraphs = []
                        for para in shape.text_frame.paragraphs:
                            runs = []
                            for run in para.runs:
                                runs.append(
                                    {
                                        "text": run.text,
                                        "font": (
                                            run.font.name
                                            if run.font and run.font.name
                                            else "DefaultFont"
                                        ),
                                        "size": (
                                            run.font.size.pt
                                            if run.font and run.font.size
                                            else None
                                        ),
                                        "bold": (
                                            run.font.bold
                                            if run.font and run.font.bold is not None
                                            else False
                                        ),
                                        "italic": (
                                            run.font.italic
                                            if run.font and run.font.italic is not None
                                            else False
                                        ),
                                        "underline": (
                                            run.font.underline
                                            if run.font
                                            and run.font.underline is not None
                                            else False
                                        ),
                                        "color": (
                                            run.font.color.rgb
                                            if run.font
                                            and run.font.color
                                            and hasattr(run.font.color, "rgb")
                                            else None
                                        ),
                                    }
                                )
                            para_text = para.text.strip()
                            if para_text:
                                paragraphs.append(
                                    {
                                        "text": para_text,
                                        "alignment": para.alignment,
                                        "runs": runs,
                                    }
                                )
                        unique_paras = []
                        seen_texts = set()
                        for para in paragraphs:
                            if para["text"] not in seen_texts:
                                unique_paras.append(para)
                                seen_texts.add(para["text"])
                        shape_data["paragraphs"] = unique_paras
                        shape_data["text"] = " ".join(
                            [p["text"] for p in unique_paras]
                        ).strip()
                    else:
                        if hasattr(shape, "text") and shape.text:
                            shape_data["text"] = shape.text.strip()
                    if "text" in shape_data and shape_data["text"]:
                        slide_items.append(shape_data)
                if slide_items:
                    text_data.append({"slide": i + 1, "shapes": slide_items})
            return text_data
        else:
            raise ValueError("Unsupported file type for text extraction")

    def extract_links(self, filepath: str):
        document = self.loader.load_file(filepath)
        links = []
        try:
            import fitz

            is_mupdf = hasattr(document, "page_count")
        except ImportError:
            is_mupdf = False
        if is_mupdf:
            seen_pdf_links = set()
            for i in range(document.page_count):
                page = document[i]
                link_list = page.get_links()
                for link in link_list:
                    uri = link.get("uri", "")
                    if uri:
                        key = (i, uri)
                        if key not in seen_pdf_links:
                            links.append(
                                {
                                    "page": i + 1,
                                    "uri": uri,
                                    "rect": list(link.get("from", [])),
                                    "kind": link.get("kind", ""),
                                }
                            )
                            seen_pdf_links.add(key)
            return links
        elif hasattr(document, "pages"):
            try:
                import pdfplumber
            except ImportError:
                pdfplumber = None
            if pdfplumber:
                try:
                    with pdfplumber.open(filepath) as pdf:
                        for i, page in enumerate(pdf.pages):
                            text = page.extract_text(layout=True)
                            found = re.findall(r"https?://\S+", text or "")
                            for url in found:
                                links.append(
                                    {"page": i + 1, "uri": url, "contents": ""}
                                )
                except Exception:
                    pass
            else:
                for i, page in enumerate(document.pages):
                    try:
                        raw_text = (
                            page.extract_text() if hasattr(page, "extract_text") else ""
                        )
                        found = re.findall(r"https?://\S+", merge_lines(raw_text))
                        for url in found:
                            links.append({"page": i + 1, "uri": url, "contents": ""})
                    except Exception:
                        pass
            return links
        elif hasattr(document, "paragraphs"):
            for para in document.paragraphs:
                if "http" in para.text:
                    links.append({"text": para.text.strip(), "url": para.text.strip()})
        elif hasattr(document, "slides"):
            seen_links = set()
            for i, slide in enumerate(document.slides):
                for shape in slide.shapes:
                    try:
                        if (
                            hasattr(shape, "click_action")
                            and shape.click_action
                            and getattr(shape.click_action, "hyperlink", None)
                        ):
                            url = str(shape.click_action.hyperlink).strip()
                            text = (
                                shape.text.strip()
                                if hasattr(shape, "text") and shape.text
                                else ""
                            )
                            if url and (url, text) not in seen_links:
                                links.append({"slide": i + 1, "text": text, "url": url})
                                seen_links.add((url, text))
                    except Exception:
                        continue
        else:
            raise ValueError("Unsupported file type for link extraction")
        return links

    def _extract_ppt_images_from_shape(self, shape, slide_index):
        images = []
        if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                images.append(
                    {
                        "slide": slide_index,
                        "image_id": shape.shape_id,
                        "blob": base64.b64encode(image.blob).decode("utf-8"),
                        "format": image.ext,
                        "width": image.width,
                        "height": image.height,
                    }
                )
            except Exception:
                pass
        elif hasattr(shape, "fill") and shape.fill.type == MSO_FILL.PICTURE:
            try:
                image = shape.fill.picture
                images.append(
                    {
                        "slide": slide_index,
                        "image_id": shape.shape_id,
                        "blob": base64.b64encode(image.blob).decode("utf-8"),
                        "format": image.ext,
                        "width": image.width,
                        "height": image.height,
                    }
                )
            except Exception:
                pass
        if hasattr(shape, "shapes"):
            for subshape in shape.shapes:
                images.extend(
                    self._extract_ppt_images_from_shape(subshape, slide_index)
                )
        else:
            try:
                blips = shape.element.xpath(".//a:blip")
                if blips:
                    rId = blips[0].get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    if (
                        rId
                        and hasattr(shape.part, "related_parts")
                        and rId in shape.part.related_parts
                    ):
                        image_part = shape.part.related_parts[rId]
                        images.append(
                            {
                                "slide": slide_index,
                                "image_id": getattr(shape, "shape_id", None),
                                "blob": base64.b64encode(image_part.blob).decode(
                                    "utf-8"
                                ),
                                "format": getattr(image_part, "ext", "unknown"),
                                "width": None,
                                "height": None,
                            }
                        )
            except Exception:
                pass
        return images

    def extract_images(self, filepath: str):
        document = self.loader.load_file(filepath)
        images = []
        try:
            import fitz

            is_mupdf = hasattr(document, "page_count")
        except ImportError:
            is_mupdf = False
        if is_mupdf:
            for i in range(document.page_count):
                page = document[i]
                image_list = page.get_images(full=True)
                for img in image_list:
                    xref = img[0]
                    base_image = document.extract_image(xref)
                    image_bytes = base_image.get("image")
                    width = base_image.get("width")
                    height = base_image.get("height")
                    image_format = base_image.get("ext")
                    if image_bytes:
                        encoded_blob = base64.b64encode(image_bytes).decode("utf-8")
                    else:
                        encoded_blob = ""
                    images.append(
                        {
                            "page": i + 1,
                            "xref": xref,
                            "format": image_format,
                            "width": width,
                            "height": height,
                            "blob": encoded_blob,
                        }
                    )
            return images
        elif hasattr(document, "pages"):
            try:
                import pdfplumber
            except ImportError:
                pdfplumber = None
            if pdfplumber:
                try:
                    with pdfplumber.open(filepath) as pdf:
                        pass
                except Exception:
                    pass
        elif hasattr(document, "inline_shapes"):
            for i, shape in enumerate(document.inline_shapes):
                try:
                    image = shape._inline.graphic.graphicData.pic.blipFill.blip
                    image_blob = shape.part.related_parts[image.embed].blob
                    encoded_blob = base64.b64encode(image_blob).decode("utf-8")
                    images.append(
                        {
                            "index": i + 1,
                            "filename": shape._inline.docPr.get("name", f"Image{i+1}"),
                            "blob": encoded_blob,
                            "width": shape.width,
                            "height": shape.height,
                        }
                    )
                except Exception:
                    continue
            if hasattr(document, "element"):
                try:
                    drawings = document.element.xpath("//w:drawing")
                    for i, drawing in enumerate(drawings):
                        try:
                            blip = drawing.xpath(".//a:blip")[0]
                            rId = blip.get(
                                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                            )
                            image_blob = document.part.related_parts[rId].blob
                            encoded_blob = base64.b64encode(image_blob).decode("utf-8")
                            images.append(
                                {
                                    "floating_index": i + 1,
                                    "rId": rId,
                                    "blob": encoded_blob,
                                }
                            )
                        except Exception:
                            continue
                except Exception:
                    pass
        elif hasattr(document, "slides"):
            for i, slide in enumerate(document.slides):
                for shape in slide.shapes:
                    extracted = self._extract_ppt_images_from_shape(shape, i + 1)
                    if extracted:
                        images.extend(extracted)
            # Fallback: If no images found via shapes, try relationship-based extraction.
            if not images:
                images = extract_images_from_rels(document)
            return images
        else:
            raise ValueError("Unsupported file type for image extraction")
        return images

    def extract_tables(self, filepath: str):
        document = self.loader.load_file(filepath)
        tables = []
        if hasattr(document, "tables"):
            for i, table in enumerate(document.tables):
                table_style = str(table.style) if table.style else "DefaultTableStyle"
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append(
                    {
                        "table_index": i + 1,
                        "style": table_style,
                        "data": table_data,
                        "rows": len(table_data),
                        "columns": (
                            len(table_data[0]) if table_data and table_data[0] else 0
                        ),
                    }
                )
        elif hasattr(document, "slides"):
            table_counter = 1
            for i, slide in enumerate(document.slides):
                for shape in slide.shapes:
                    try:
                        if shape.has_table:
                            table_data = []
                            table = shape.table
                            for row in table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(row_data)
                            if table_data:
                                tables.append(
                                    {
                                        "table_index": table_counter,
                                        "style": "DefaultPPTTable",
                                        "slide": i + 1,
                                        "data": table_data,
                                        "rows": len(table_data),
                                        "columns": (
                                            len(table_data[0])
                                            if table_data and table_data[0]
                                            else 0
                                        ),
                                    }
                                )
                                table_counter += 1
                    except Exception:
                        continue
        elif hasattr(document, "page_count"):
            try:
                import pdfplumber
            except ImportError:
                pdfplumber = None
            if pdfplumber:
                try:
                    with pdfplumber.open(filepath) as pdf:
                        for i, page in enumerate(pdf.pages):
                            table_data = page.extract_table()
                            if table_data:
                                tables.append(
                                    {
                                        "page": i + 1,
                                        "style": "N/A",
                                        "data": table_data,
                                        "rows": len(table_data),
                                        "columns": (
                                            len(table_data[0])
                                            if table_data and table_data[0]
                                            else 0
                                        ),
                                    }
                                )
                    if not tables:
                        tables.append({"warning": "No tables found with pdfplumber."})
                except Exception as e:
                    tables.append(
                        {
                            "warning": f"Error extracting tables with pdfplumber: {str(e)}"
                        }
                    )
            else:
                tables.append(
                    {
                        "warning": "pdfplumber not installed. Install pdfplumber for PDF table extraction."
                    }
                )
        else:
            raise ValueError("Unsupported file type for table extraction")
        return tables
